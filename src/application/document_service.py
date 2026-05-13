"""
Application Layer - Document Service

Use cases for document ingestion and management.
Supports multiple PDF backends for flexible extraction.
"""

from __future__ import annotations

import hashlib
import inspect
import logging
import re
import shutil
import tempfile
import time
from collections.abc import Awaitable, Callable
from pathlib import Path
from typing import TYPE_CHECKING, Any, cast

from src.application.citation_artifacts import (
    empty_citation_reason,
    remove_citation_index,
    save_citation_status,
)
from src.application.markdown_block_builder import build_markdown_blocks
from src.application.output_paths import resolve_document_output_path
from src.domain.citation import build_evidence_spans
from src.domain.entities import (
    DocumentManifest,
    DocumentSummary,
    FigureAsset,
    IngestResult,
    SectionAsset,
    TableAsset,
)
from src.domain.etl_profile import ETLProfile
from src.domain.line_spans import (
    MarkdownLineSpanIndex,
    annotate_marker_blocks,
    apply_asset_line_spans,
)
from src.domain.marker_errors import (
    format_marker_failure,
    is_marker_backend_unavailable,
    is_marker_resource_error,
)
from src.domain.services import ManifestGenerator
from src.domain.value_objects import DocId

if TYPE_CHECKING:
    from src.domain.repositories import (
        DocumentRepository,
        KnowledgeGraphInterface,
        PDFExtractorInterface,
    )
    from src.infrastructure.marker_adapter import MarkerPDFExtractor
    from src.infrastructure.ocr_processor import OCRProcessor


ToolProgressCallback = Callable[[int, int, str, str], Awaitable[None] | None]
PageRange = tuple[int, int]

_PAGE_MARKER_RE = re.compile(r"<!-- Page (\d+) -->")
logger = logging.getLogger(__name__)


def format_page_ranges(page_ranges: list[PageRange] | tuple[PageRange, ...]) -> str:
    """Format normalized inclusive page ranges for logs and doc_id scopes."""
    parts = []
    for start_page, end_page in page_ranges:
        if start_page == end_page:
            parts.append(str(start_page))
        else:
            parts.append(f"{start_page}-{end_page}")
    return ",".join(parts)


def build_page_number_map(
    page_ranges: list[PageRange] | tuple[PageRange, ...],
) -> list[int]:
    """Expand normalized ranges into a sequential subset→original page map."""
    page_numbers: list[int] = []
    for start_page, end_page in page_ranges:
        page_numbers.extend(range(start_page, end_page + 1))
    return page_numbers


def normalize_page_ranges(
    page_ranges: list[str] | None,
    total_pages: int,
) -> tuple[PageRange, ...]:
    """Validate and merge user-supplied 1-indexed inclusive page ranges."""
    if not page_ranges:
        return ()

    normalized: list[PageRange] = []
    for raw_spec in page_ranges:
        spec = raw_spec.strip()
        if not spec:
            continue

        if "-" in spec:
            start_text, end_text = spec.split("-", 1)
            start_page = int(start_text)
            end_page = int(end_text)
        else:
            start_page = int(spec)
            end_page = start_page

        if start_page < 1 or end_page < 1:
            raise ValueError("Page numbers must be >= 1")
        if start_page > end_page:
            raise ValueError(f"Invalid page range: {spec}")
        if end_page > total_pages:
            raise ValueError(
                f"Page range {spec} exceeds total page count {total_pages}"
            )

        normalized.append((start_page, end_page))

    if not normalized:
        return ()

    normalized.sort()
    merged: list[PageRange] = [normalized[0]]
    for start_page, end_page in normalized[1:]:
        prev_start, prev_end = merged[-1]
        if start_page <= prev_end + 1:
            merged[-1] = (prev_start, max(prev_end, end_page))
        else:
            merged.append((start_page, end_page))
    return tuple(merged)


def remap_page_number(page_number: int, page_map: list[int] | None) -> int:
    """Translate subset-local page numbers back to original PDF page numbers."""
    if not page_map or page_number < 1 or page_number > len(page_map):
        return page_number
    return page_map[page_number - 1]


def build_doc_id_unique_suffix(
    source_path: Path,
    page_ranges: list[PageRange] | tuple[PageRange, ...] | None = None,
) -> str:
    """Build a stable DocId uniqueness suffix that includes page scoping."""
    suffix = str(source_path.absolute())
    if page_ranges:
        suffix = f"{suffix}#pages={format_page_ranges(page_ranges)}"
    return suffix


def materialize_pdf_page_subset(
    source_path: Path,
    output_path: Path,
    page_ranges: list[PageRange] | tuple[PageRange, ...],
) -> Path:
    """Persist a subset PDF containing only the requested inclusive page ranges."""
    import fitz  # type: ignore

    output_path.parent.mkdir(parents=True, exist_ok=True)
    subset_pdf = fitz.open()
    try:
        with fitz.open(str(source_path)) as source_pdf:
            for start_page, end_page in page_ranges:
                subset_pdf.insert_pdf(
                    source_pdf,
                    from_page=start_page - 1,
                    to_page=end_page - 1,
                )
        subset_pdf.save(output_path)
    finally:
        subset_pdf.close()
    return output_path


def remap_markdown_page_markers(markdown: str, page_map: list[int] | None) -> str:
    """Rewrite subset-local markdown page markers to original PDF numbers."""
    if not page_map:
        return markdown

    marker_index = 0

    def replace_page_marker(match: re.Match[str]) -> str:
        nonlocal marker_index
        if marker_index >= len(page_map):
            return match.group(0)
        original_page = page_map[marker_index]
        marker_index += 1
        return f"<!-- Page {original_page} -->"

    return _PAGE_MARKER_RE.sub(replace_page_marker, markdown)


def remap_toc_pages(
    toc: list[tuple[int, str, int]],
    page_map: list[int] | None,
) -> list[tuple[int, str, int]]:
    """Translate PDF TOC page numbers from subset-local to original numbering."""
    if not page_map:
        return toc
    return [
        (level, title, remap_page_number(page_number, page_map))
        for level, title, page_number in toc
    ]


async def _invoke_progress_callback(
    callback: ToolProgressCallback | None,
    step: int,
    total_steps: int,
    phase: str,
    message: str,
) -> None:
    if callback is None:
        return

    result = callback(step, total_steps, phase, message)
    if inspect.isawaitable(result):
        await result


class DocumentService:
    """
    Application service for document operations.

    Orchestrates the ETL pipeline:
    1. Extract text and images from PDF
    2. Generate document manifest
    3. Index in knowledge graph (optional)
    4. Save to repository

    Supports multiple PDF backends:
    - PyMuPDF (default, fast, no models)
    - Marker (optional, use_marker=True, produces blocks.json with bbox/section_hierarchy)
    """

    def __init__(
        self,
        repository: DocumentRepository,
        pdf_extractor: PDFExtractorInterface,
        knowledge_graph: KnowledgeGraphInterface | None = None,
        marker_extractor: MarkerPDFExtractor | None = None,
        profile: ETLProfile | None = None,
        ocr_processor: OCRProcessor | None = None,
    ):
        """
        Initialize document service with dependencies.

        Args:
            repository: Document storage repository
            pdf_extractor: PDF extraction implementation (PyMuPDF)
            knowledge_graph: Optional knowledge graph for indexing
            marker_extractor: Optional Marker PDF extractor for structured parsing
            profile: Optional ETL profile (auto-detected from pdf_extractor if not provided)
        """
        self.repository = repository
        self.pdf_extractor = pdf_extractor
        self.knowledge_graph = knowledge_graph
        self.marker_extractor = marker_extractor
        self.ocr_processor = ocr_processor

        # Resolve profile: explicit > from extractor > default
        extractor_profile = getattr(pdf_extractor, "profile", None)
        if profile is not None:
            resolved_profile = profile
        elif isinstance(extractor_profile, ETLProfile):
            resolved_profile = extractor_profile
        else:
            resolved_profile = ETLProfile.default()

        self.profile = resolved_profile
        self.manifest_generator = ManifestGenerator(profile=resolved_profile)

    async def ingest(
        self,
        file_paths: list[str],
        use_marker: bool = False,
        progress_callback: ToolProgressCallback | None = None,
        *,
        ocr_enabled: bool = False,
        ocr_language: str = "eng",
        rotate_pages: bool = False,
        deskew: bool = False,
        marker_max_pages_per_chunk: int = 0,
        extract_figures: bool = True,
        page_ranges: list[str] | None = None,
        require_marker: bool = False,
    ) -> list[IngestResult]:
        """
        Ingest multiple PDF files.

        Args:
            file_paths: List of paths to PDF files
            use_marker: If True, use Marker for structured parsing (slower but richer)
                        - Produces blocks.json with bbox and section_hierarchy
                        - Better TOC and figure caption extraction
                        - Requires Marker models (~1GB first run)

        Returns:
            List of IngestResult for each file
        """
        results = []
        base_steps = 9 if use_marker and self.marker_extractor else 8
        per_file_steps = base_steps + (1 if ocr_enabled else 0)
        total_steps = max(len(file_paths), 1) * per_file_steps

        for index, file_path in enumerate(file_paths):
            base_step = index * per_file_steps

            async def file_progress(
                step: int,
                inner_total_steps: int,
                phase: str,
                message: str,
                *,
                _base_step: int = base_step,
                _file_index: int = index,
            ) -> None:
                bounded_total = max(inner_total_steps, 1)
                mapped_step = _base_step + min(max(step, 0), bounded_total)
                await _invoke_progress_callback(
                    progress_callback,
                    mapped_step,
                    total_steps,
                    phase,
                    f"[{_file_index + 1}/{len(file_paths)}] {message}",
                )

            if use_marker and self.marker_extractor:
                result = await self._ingest_single_with_marker(
                    file_path,
                    progress_callback=file_progress,
                    ocr_enabled=ocr_enabled,
                    ocr_language=ocr_language,
                    rotate_pages=rotate_pages,
                    deskew=deskew,
                    marker_max_pages_per_chunk=marker_max_pages_per_chunk,
                    extract_figures=extract_figures,
                    page_ranges=page_ranges,
                    require_marker=require_marker,
                )
            else:
                if use_marker and require_marker:
                    path = Path(file_path)
                    results.append(
                        IngestResult(
                            doc_id="",
                            filename=path.name,
                            success=False,
                            error=(
                                "Marker structure parse was required, but the "
                                "Marker extractor is temporarily unavailable because "
                                "marker-pdf 1.10.2 pins Pillow<11 while the secure "
                                "runtime requires Pillow>=12.2.0."
                            ),
                        )
                    )
                    continue
                result = await self._ingest_single(
                    file_path,
                    progress_callback=file_progress,
                    ocr_enabled=ocr_enabled,
                    ocr_language=ocr_language,
                    rotate_pages=rotate_pages,
                    deskew=deskew,
                    page_ranges=page_ranges,
                )
                if use_marker and result.success:
                    result.backend = "pymupdf_fallback"
                    result.warnings.append(
                        "Marker was requested but the extractor was not configured; "
                        "used PyMuPDF fallback while marker-pdf depends on vulnerable "
                        "Pillow<11."
                    )
            results.append(result)

        return results

    async def _ingest_single(
        self,
        file_path: str,
        progress_callback: ToolProgressCallback | None = None,
        *,
        ocr_enabled: bool = False,
        ocr_language: str = "eng",
        rotate_pages: bool = False,
        deskew: bool = False,
        page_ranges: list[str] | None = None,
    ) -> IngestResult:
        """Ingest a single PDF file."""
        start_time = time.time()
        path = Path(file_path)
        total_steps = 9 if ocr_enabled else 8

        # Validate file exists
        if not path.exists():
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=f"File not found: {path}",
            )

        if path.suffix.lower() != ".pdf":
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=f"Not a PDF file: {path}",
            )

        # Validate PDF magic bytes (%PDF-)
        try:
            with path.open("rb") as f:
                header = f.read(5)
            if header != b"%PDF-":
                return IngestResult(
                    doc_id="",
                    filename=path.name,
                    success=False,
                    error="Invalid PDF: file does not start with %PDF- header",
                )
        except OSError as e:
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=f"Cannot read file: {e}",
            )

        try:
            total_page_count = self.pdf_extractor.get_page_count(path)
            source_pdf_sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
            normalized_page_ranges = normalize_page_ranges(
                page_ranges, total_page_count
            )
            page_map = build_page_number_map(normalized_page_ranges)
            processed_page_count = len(page_map) if page_map else total_page_count
            reported_page_count = total_page_count if page_map else processed_page_count

            # Generate unique doc_id
            await _invoke_progress_callback(
                progress_callback,
                1,
                total_steps,
                "Preparing",
                f"Preparing {path.name}",
            )
            doc_id = DocId.generate(
                path.stem,
                build_doc_id_unique_suffix(path, normalized_page_ranges),
            )
            self._save_original_pdf_copy(doc_id.value, path)

            active_pdf_path = path
            if normalized_page_ranges:
                active_pdf_path = materialize_pdf_page_subset(
                    path,
                    self.repository.get_doc_dir(doc_id.value) / "selected_pages.pdf",
                    normalized_page_ranges,
                )

            current_step = 2
            if ocr_enabled:
                await _invoke_progress_callback(
                    progress_callback,
                    current_step,
                    total_steps,
                    "OCR Preprocessing",
                    f"Running OCR preprocessing for {path.name}",
                )
                active_pdf_path = self._preprocess_pdf_with_ocr(
                    doc_id.value,
                    active_pdf_path,
                    language=ocr_language,
                    rotate_pages=rotate_pages,
                    deskew=deskew,
                )
                current_step += 1

            # Step 1: Extract text as markdown
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Extracting Text",
                f"Extracting text from {path.name}",
            )
            markdown = self.pdf_extractor.extract_text(active_pdf_path)
            markdown = remap_markdown_page_markers(markdown, page_map)
            current_step += 1

            # Step 2: Save markdown
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Saving Markdown",
                f"Saving markdown for {path.name}",
            )
            markdown_path = self.repository.save_markdown(doc_id.value, markdown)
            current_step += 1

            # Step 3: Extract and save images
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Extracting Figures",
                f"Extracting figures from {path.name}",
            )
            figures = await self._extract_and_save_images(
                doc_id.value,
                active_pdf_path,
                page_map=page_map or None,
            )
            current_step += 1

            # Step 3.5: Extract tables (Docling enhanced)
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Extracting Tables",
                f"Extracting tables from {path.name}",
            )
            tables = await self._extract_tables(
                active_pdf_path,
                page_map=page_map or None,
            )
            apply_asset_line_spans(MarkdownLineSpanIndex(markdown), figures, tables)
            current_step += 1

            # Step 4: Get page count
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Reading Metadata",
                f"Reading PDF metadata from {path.name}",
            )
            page_count = reported_page_count
            current_step += 1

            # Step 4.5: Get PDF built-in TOC and metadata title (if available)
            pdf_toc: list[tuple[int, str, int]] = []
            pdf_title = ""
            if hasattr(self.pdf_extractor, "get_toc"):
                pdf_toc = remap_toc_pages(
                    self.pdf_extractor.get_toc(active_pdf_path),
                    page_map or None,
                )
            if hasattr(self.pdf_extractor, "get_title"):
                pdf_title = self.pdf_extractor.get_title(active_pdf_path)

            # Step 5: Extract entities from knowledge graph (if available)
            entities = []
            if self.knowledge_graph and self.knowledge_graph.is_available:
                await _invoke_progress_callback(
                    progress_callback,
                    current_step,
                    total_steps,
                    "Indexing Knowledge Graph",
                    f"Indexing {path.name} into the knowledge graph",
                )
                try:
                    # Index the document
                    await self.knowledge_graph.insert(doc_id.value, markdown)
                    # Extract entities
                    entities = await self.knowledge_graph.extract_entities(markdown)
                except Exception:
                    # Log but don't fail - LightRAG is optional
                    import logging

                    logging.getLogger(__name__).warning(
                        "LightRAG indexing failed for %s", doc_id.value, exc_info=True
                    )
            else:
                await _invoke_progress_callback(
                    progress_callback,
                    current_step,
                    total_steps,
                    "Skipping Knowledge Graph",
                    f"Knowledge graph indexing skipped for {path.name}",
                )
            current_step += 1

            # Step 6: Generate manifest
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Generating Manifest",
                f"Generating manifest for {path.name}",
            )
            manifest = self.manifest_generator.generate(
                doc_id=doc_id.value,
                filename=path.name,
                markdown=markdown,
                figures=figures,
                tables=tables,  # Pass Docling-extracted tables
                page_count=page_count,
                markdown_path=str(markdown_path),
                lightrag_entities=entities,
                pdf_toc=pdf_toc,
                pdf_title=pdf_title,
                source_pdf_sha256=source_pdf_sha256,
                selected_page_map=page_map,
            )

            # Step 7: Save manifest
            self.repository.save_manifest(manifest)
            blocks_data = build_markdown_blocks(markdown, manifest)
            segmentation_warnings: list[str] = []
            if blocks_data:
                self._save_blocks_json(doc_id.value, blocks_data)
                self._save_citation_index(
                    doc_id.value,
                    markdown,
                    blocks_data,
                    source_backend="pymupdf",
                )
                await self._save_segmentation_artifact(
                    doc_id.value,
                    segmentation_warnings,
                )
            await _invoke_progress_callback(
                progress_callback,
                total_steps,
                total_steps,
                "Completed",
                f"Finished processing {path.name}",
            )

            processing_time = time.time() - start_time

            return IngestResult(
                doc_id=doc_id.value,
                filename=path.name,
                title=manifest.title,
                success=True,
                manifest=manifest,
                pages_processed=processed_page_count,
                tables_found=len(manifest.assets.tables),
                figures_found=len(manifest.assets.figures),
                sections_found=len(manifest.assets.sections),
                processing_time_seconds=processing_time,
                warnings=segmentation_warnings,
            )

        except Exception as e:
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=str(e),
            )

    async def _ingest_single_with_marker(
        self,
        file_path: str,
        progress_callback: ToolProgressCallback | None = None,
        *,
        ocr_enabled: bool = False,
        ocr_language: str = "eng",
        rotate_pages: bool = False,
        deskew: bool = False,
        marker_max_pages_per_chunk: int = 0,
        extract_figures: bool = True,
        page_ranges: list[str] | None = None,
        require_marker: bool = False,
    ) -> IngestResult:
        """
        Ingest a single PDF file using Marker for structured parsing.

        This provides richer structure than PyMuPDF:
        - blocks.json with bbox and section_hierarchy
        - Better TOC extraction
        - Figure caption association
        """
        start_time = time.time()
        path = Path(file_path)
        total_steps = 10 if ocr_enabled else 9

        # Validate file exists
        if not path.exists():
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=f"File not found: {path}",
            )

        if path.suffix.lower() != ".pdf":
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=f"Not a PDF file: {path}",
            )

        if self.marker_extractor is None:
            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error="Marker extractor not available",
            )

        try:
            total_page_count = self.pdf_extractor.get_page_count(path)
            source_pdf_sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
            normalized_page_ranges = normalize_page_ranges(
                page_ranges, total_page_count
            )
            page_map = build_page_number_map(normalized_page_ranges)
            processed_page_count = len(page_map) if page_map else total_page_count
            reported_page_count = total_page_count if page_map else processed_page_count

            # Generate unique doc_id
            await _invoke_progress_callback(
                progress_callback,
                1,
                total_steps,
                "Preparing",
                f"Preparing {path.name}",
            )
            doc_id = DocId.generate(
                path.stem,
                build_doc_id_unique_suffix(path, normalized_page_ranges),
            )
            self._save_original_pdf_copy(doc_id.value, path)

            active_pdf_path = path
            if normalized_page_ranges:
                active_pdf_path = materialize_pdf_page_subset(
                    path,
                    self.repository.get_doc_dir(doc_id.value) / "selected_pages.pdf",
                    normalized_page_ranges,
                )

            current_step = 2
            if ocr_enabled:
                await _invoke_progress_callback(
                    progress_callback,
                    current_step,
                    total_steps,
                    "OCR Preprocessing",
                    f"Running OCR preprocessing for {path.name}",
                )
                active_pdf_path = self._preprocess_pdf_with_ocr(
                    doc_id.value,
                    active_pdf_path,
                    language=ocr_language,
                    rotate_pages=rotate_pages,
                    deskew=deskew,
                )
                current_step += 1

            # Step 1: Parse PDF with Marker (rich structure)
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Parsing Structure",
                f"Loading Marker models and parsing structure from {path.name}",
            )
            parse_result = self.marker_extractor.parse(
                active_pdf_path,
                extract_images=extract_figures,
                max_pages_per_chunk=(
                    marker_max_pages_per_chunk
                    if marker_max_pages_per_chunk > 0
                    else None
                ),
                page_map=page_map or None,
                reported_page_count=reported_page_count if page_map else None,
            )
            current_step += 1
            warnings: list[str] = []
            if not str(parse_result.markdown or "").strip():
                return IngestResult(
                    doc_id=doc_id.value,
                    filename=path.name,
                    success=False,
                    error=(
                        "Marker returned empty markdown; retry with OCR enabled "
                        "or use the PyMuPDF backend."
                    ),
                    backend="marker",
                    warnings=[
                        "Marker returned empty markdown before citation artifacts "
                        "could be created."
                    ],
                )

            # Step 2: Save markdown
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Saving Markdown",
                f"Saving markdown for {path.name}",
            )
            line_span_index = annotate_marker_blocks(
                parse_result.markdown,
                parse_result.blocks,
            )
            markdown_path = self.repository.save_markdown(
                doc_id.value, parse_result.markdown
            )
            current_step += 1

            # Step 3: Prepare blocks.json (structured data)
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Preparing Blocks",
                f"Preparing structured blocks for {path.name}",
            )
            blocks_data = self._convert_blocks_to_json(parse_result.blocks)
            current_step += 1

            # Step 4: Extract and save images from Marker result
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Extracting Figures",
                f"Extracting figures from {path.name}",
            )
            figures = await self._save_marker_images(
                doc_id.value,
                parse_result,
                pdf_path=path,
            )
            current_step += 1

            # Step 5: Convert Marker blocks to TableAsset
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Extracting Tables",
                f"Extracting tables from {path.name}",
            )
            tables = self._extract_tables_from_blocks(parse_result.blocks)
            current_step += 1

            # Step 6: Convert TOC to SectionAsset
            await _invoke_progress_callback(
                progress_callback,
                current_step,
                total_steps,
                "Building Sections",
                f"Building section tree for {path.name}",
            )
            sections = self._extract_sections_from_toc(
                parse_result.toc,
                parse_result.markdown,
            )
            apply_asset_line_spans(
                line_span_index,
                figures,
                tables,
                blocks=parse_result.blocks,
                sections=sections,
            )
            current_step += 1

            # Step 7: Get page count
            page_count = parse_result.page_count or reported_page_count

            # Step 8: Index in knowledge graph (if available)
            entities = []
            if self.knowledge_graph and self.knowledge_graph.is_available:
                await _invoke_progress_callback(
                    progress_callback,
                    current_step,
                    total_steps,
                    "Indexing Knowledge Graph",
                    f"Indexing {path.name} into the knowledge graph",
                )
                try:
                    await self.knowledge_graph.insert(
                        doc_id.value, parse_result.markdown
                    )
                    entities = await self.knowledge_graph.extract_entities(
                        parse_result.markdown
                    )
                except Exception as e:
                    import logging

                    logging.warning(f"LightRAG indexing failed: {e}")
            else:
                await _invoke_progress_callback(
                    progress_callback,
                    current_step,
                    total_steps,
                    "Skipping Knowledge Graph",
                    f"Knowledge graph indexing skipped for {path.name}",
                )
            current_step += 1

            # Step 9: Generate manifest (with richer data)
            # Note: sections are parsed from markdown by ManifestGenerator
            manifest = self.manifest_generator.generate(
                doc_id=doc_id.value,
                filename=path.name,
                markdown=parse_result.markdown,
                figures=figures,
                tables=tables,
                page_count=page_count,
                markdown_path=str(markdown_path),
                lightrag_entities=entities,
                sections=sections,
                source_pdf_sha256=source_pdf_sha256,
                selected_page_map=page_map,
            )

            # Step 10: Save manifest
            citation_backend = "marker"
            if self._marker_blocks_need_markdown_fallback(blocks_data, manifest):
                citation_backend = "marker_markdown_fallback"
                blocks_data = build_markdown_blocks(
                    parse_result.markdown,
                    manifest,
                    source_backend=citation_backend,
                )
                warnings.append(
                    "Marker emitted non-citeable layout blocks; synthesized "
                    "markdown-based blocks for citation and segmentation."
                )
            self.repository.save_manifest(manifest)
            self._save_blocks_json(doc_id.value, blocks_data)
            self._save_citation_index(
                doc_id.value,
                parse_result.markdown,
                blocks_data,
                source_backend=citation_backend,
            )
            await self._save_segmentation_artifact(doc_id.value, warnings)
            await _invoke_progress_callback(
                progress_callback,
                total_steps,
                total_steps,
                "Completed",
                f"Finished processing {path.name}",
            )

            processing_time = time.time() - start_time

            return IngestResult(
                doc_id=doc_id.value,
                filename=path.name,
                title=manifest.title or parse_result.metadata.get("title", ""),
                success=True,
                manifest=manifest,
                pages_processed=processed_page_count,
                tables_found=len(manifest.assets.tables),
                figures_found=len(manifest.assets.figures),
                sections_found=len(manifest.assets.sections),
                processing_time_seconds=processing_time,
                backend="marker",  # Indicate which backend was used
                warnings=warnings,
            )

        except Exception as e:
            marker_message = format_marker_failure(e)
            if is_marker_backend_unavailable(e) or is_marker_resource_error(e):
                if require_marker:
                    return IngestResult(
                        doc_id="",
                        filename=path.name,
                        success=False,
                        error=marker_message,
                    )
                fallback_result = await self._ingest_single(
                    file_path,
                    progress_callback=progress_callback,
                    ocr_enabled=ocr_enabled,
                    ocr_language=ocr_language,
                    rotate_pages=rotate_pages,
                    deskew=deskew,
                    page_ranges=page_ranges,
                )
                if fallback_result.success:
                    fallback_result.backend = "pymupdf_fallback"
                    fallback_result.warnings.append(
                        f"Marker parse failed and PyMuPDF fallback was used. {marker_message}"
                    )
                    return fallback_result

                return IngestResult(
                    doc_id="",
                    filename=path.name,
                    success=False,
                    error=(
                        f"{marker_message}\n"
                        f"PyMuPDF fallback also failed: {fallback_result.error}"
                    ),
                )

            return IngestResult(
                doc_id="",
                filename=path.name,
                success=False,
                error=marker_message,
            )

    def _convert_blocks_to_json(self, blocks: list) -> list[dict]:
        """Convert MarkerBlock objects to JSON-serializable dicts."""
        return [
            {
                "block_id": b.block_id,
                "block_type": b.block_type,
                "page": b.page,
                "text": b.text or "",
                "text_preview": (b.text[:500] if b.text else ""),
                "bbox": b.bbox,
                "polygon": b.polygon,
                "section_hierarchy": b.section_hierarchy,
                "metadata": b.metadata,
            }
            for b in blocks
        ]

    def _save_blocks_json(self, doc_id: str, blocks_data: list[dict]) -> Path:
        """Save blocks.json to repository."""
        return self.repository.save_blocks(doc_id, blocks_data)

    def _save_citation_index(
        self,
        doc_id: str,
        markdown: str,
        blocks_data: list[dict[str, Any]],
        *,
        source_backend: str,
    ) -> Path | None:
        """Build and persist citation-ready spans for downstream references."""
        spans = build_evidence_spans(
            doc_id=doc_id,
            markdown=markdown,
            blocks=blocks_data,
            source_backend=source_backend,
        )
        try:
            save_citation_status(
                self.repository,
                doc_id,
                source_backend=source_backend,
                found=len(spans),
                reason="" if spans else empty_citation_reason(blocks_data),
            )
        except Exception:
            logger.warning(
                "Failed to save citation extraction status for %s",
                doc_id,
                exc_info=True,
            )
        if not spans:
            try:
                remove_citation_index(self.repository, doc_id)
            except Exception:
                logger.warning(
                    "Failed to remove empty citation index for %s",
                    doc_id,
                    exc_info=True,
                )
            return None
        return self.repository.save_citation_index(doc_id, spans)

    @staticmethod
    def _marker_blocks_need_markdown_fallback(
        blocks_data: list[dict[str, Any]],
        manifest: DocumentManifest,
    ) -> bool:
        if not blocks_data:
            return True
        meaningful_types = {
            str(block.get("block_type") or "").lower() for block in blocks_data
        }
        has_text = any(str(block.get("text") or "").strip() for block in blocks_data)
        if not has_text:
            return True
        if manifest.assets.sections and "sectionheader" not in meaningful_types:
            return True
        if manifest.assets.tables and "table" not in meaningful_types:
            return True
        return meaningful_types <= {"markdownoutput"}

    async def _save_segmentation_artifact(
        self,
        doc_id: str,
        warnings: list[str],
    ) -> None:
        try:
            from src.application.segmentation_service import SegmentationService

            await SegmentationService(self.repository).save_document_segmentation(
                doc_id
            )
        except Exception as exc:
            logger.warning(
                "Failed to save segmentation artifact for %s",
                doc_id,
                exc_info=True,
            )
            warnings.append(f"Segmentation export skipped: {exc}")

    def _save_original_pdf_copy(self, doc_id: str, source_path: Path) -> None:
        """Persist the original PDF for overlay inspection and downstream tooling."""
        doc_dir = self.repository.get_doc_dir(doc_id)
        doc_dir.mkdir(parents=True, exist_ok=True)
        target = doc_dir / "original.pdf"
        shutil.copy2(source_path, target)

    def _preprocess_pdf_with_ocr(
        self,
        doc_id: str,
        source_path: Path,
        *,
        language: str,
        rotate_pages: bool,
        deskew: bool,
    ) -> Path:
        if self.ocr_processor is None:
            raise RuntimeError(
                "OCR preprocessing requested, but OCR processor is not configured."
            )

        doc_dir = self.repository.get_doc_dir(doc_id)
        target = doc_dir / "ocr_processed.pdf"
        result = self.ocr_processor.preprocess_pdf(
            source_path,
            target,
            language=language,
            rotate_pages=rotate_pages,
            deskew=deskew,
        )
        return result.output_path

    @staticmethod
    def _get_image_dimensions(img_bytes: bytes) -> tuple[int, int]:
        """Read image dimensions from bytes using PIL."""
        try:
            import io

            from PIL import Image

            img = Image.open(io.BytesIO(img_bytes))
            size: tuple[int, int] = img.size  # (width, height)
            return size
        except Exception:  # PIL can raise various errors
            return (0, 0)

    @staticmethod
    def _image_variance(img_bytes: bytes) -> float:
        """Estimate whether an image has enough visual variation to be useful."""
        try:
            import io

            from PIL import Image, ImageStat

            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            stat = ImageStat.Stat(img)
            return float(sum(stat.var) / len(stat.var))
        except Exception:
            return 0.0

    @staticmethod
    def _normalize_marker_image_key(value: str) -> str:
        """Normalize a Marker image filename or block path to a comparable key."""
        stem = Path(value).stem if "." in value else value
        return stem if not stem.startswith("/") else stem.replace("/", "_")

    @staticmethod
    def _get_marker_image_blocks(blocks: list) -> list:
        """Prefer semantic Figure blocks and only fall back to Picture blocks."""
        figure_blocks = [block for block in blocks if block.block_type == "Figure"]
        if figure_blocks:
            return figure_blocks
        return [block for block in blocks if block.block_type in {"Figure", "Picture"}]

    def _get_min_figure_px(self) -> int:
        """Return the configured minimum figure dimension threshold."""
        return self.profile.filters.min_figure_px

    @staticmethod
    def _render_pdf_block_image(
        pdf_path: Path,
        *,
        page_number: int,
        bbox: list[float],
        zoom: float = 2.0,
        padding: float = 6.0,
    ) -> bytes | None:
        """Render a PDF clip from a semantic Marker block bbox."""
        if len(bbox) != 4 or page_number < 1:
            return None

        try:
            import fitz
        except Exception:
            return None

        try:
            with fitz.open(str(pdf_path)) as doc:
                if page_number > len(doc):
                    return None
                page = doc[page_number - 1]
                clip = fitz.Rect(*bbox)
                if clip.is_empty:
                    return None
                clip = fitz.Rect(
                    clip.x0 - padding,
                    clip.y0 - padding,
                    clip.x1 + padding,
                    clip.y1 + padding,
                )
                clip = clip & page.rect
                if clip.is_empty:
                    return None
                pix = page.get_pixmap(
                    matrix=fitz.Matrix(zoom, zoom),
                    clip=clip,
                    alpha=False,
                )
                return cast("bytes", pix.tobytes("png"))
        except Exception:
            return None

    async def _save_marker_images(
        self, doc_id: str, parse_result: Any, *, pdf_path: Path | None = None
    ) -> list[FigureAsset]:
        """Save images from Marker parse result."""
        figures: list[FigureAsset] = []

        image_blocks = self._get_marker_image_blocks(parse_result.blocks)
        min_px = self._get_min_figure_px()
        saved_count = 0

        for matched_block in image_blocks:
            block_key = self._normalize_marker_image_key(
                str(matched_block.metadata.get("id") or "")
            )
            img_bytes = parse_result.images.get(block_key) or next(
                (
                    payload
                    for name, payload in parse_result.images.items()
                    if self._normalize_marker_image_key(name) == block_key
                ),
                None,
            )
            ext = "png"
            if img_bytes is None and pdf_path is not None:
                img_bytes = self._render_pdf_block_image(
                    pdf_path,
                    page_number=matched_block.page,
                    bbox=list(matched_block.bbox),
                )
            else:
                for name in parse_result.images:
                    if (
                        self._normalize_marker_image_key(name) == block_key
                        and "." in name
                    ):
                        ext = name.split(".")[-1]
                        break

            if img_bytes is None:
                continue

            width, height = self._get_image_dimensions(img_bytes)
            if width < min_px or height < min_px:
                continue

            saved_count += 1
            fig_id = f"fig_{matched_block.page}_{saved_count}"

            # Save image
            image_path = self.repository.save_image(
                doc_id=doc_id,
                image_id=fig_id,
                data=img_bytes,
                ext=ext,
            )

            figures.append(
                FigureAsset(
                    id=fig_id,
                    page=matched_block.page,
                    path=str(image_path),
                    ext=ext,
                    width=width,
                    height=height,
                    caption=str(matched_block.metadata.get("caption") or ""),
                    figure_type="",
                    source="marker",
                    source_block_id=matched_block.block_id,
                    source_order=int(matched_block.metadata.get("source_order") or 0),
                    line_start=int(matched_block.metadata.get("line_start"))
                    if isinstance(matched_block.metadata.get("line_start"), int)
                    else None,
                    line_end=int(matched_block.metadata.get("line_end"))
                    if isinstance(matched_block.metadata.get("line_end"), int)
                    else None,
                    line_source=str(
                        matched_block.metadata.get("line_match_strategy") or ""
                    ),
                    raw_path="",
                    figure_bbox=self._bbox_list(matched_block.bbox),
                    crop_bbox=self._bbox_list(matched_block.bbox),
                    caption_bbox=self._bbox_list(
                        matched_block.metadata.get("caption_bbox")
                    ),
                    caption_confidence=1.0
                    if str(matched_block.metadata.get("caption") or "").strip()
                    else 0.0,
                    extraction_strategy="marker_block_render",
                )
            )

        return figures

    @staticmethod
    def _parse_table_dimensions(markdown: str) -> tuple[int, int]:
        """Parse row_count and col_count from markdown table text."""
        if not markdown:
            return (0, 0)
        lines = [line.strip() for line in markdown.strip().splitlines() if line.strip()]
        # Filter out separator lines like |---|---|
        data_lines = [line for line in lines if not all(c in "-| :" for c in line)]
        row_count = len(data_lines)
        col_count = 0
        if data_lines:
            # Count columns from first data line
            col_count = data_lines[0].count("|") - 1
            if col_count < 0:
                col_count = 0
        return (row_count, col_count)

    def _extract_tables_from_blocks(self, blocks: list) -> list[TableAsset]:
        """Extract tables from Marker blocks."""
        tables = []
        table_idx = 0

        for block in blocks:
            if block.block_type == "Table":
                table_idx += 1
                row_count, col_count = self._parse_table_dimensions(block.text)
                tables.append(
                    TableAsset(
                        id=f"tab_{table_idx}",
                        page=block.page,
                        caption="",
                        preview=block.text[:100] if block.text else "",
                        markdown=block.text or "",
                        row_count=row_count,
                        col_count=col_count,
                        has_header=True,
                        source="marker",
                        source_block_id=block.block_id,
                        source_order=int(block.metadata.get("source_order") or 0),
                        line_start=int(block.metadata.get("line_start"))
                        if isinstance(block.metadata.get("line_start"), int)
                        else None,
                        line_end=int(block.metadata.get("line_end"))
                        if isinstance(block.metadata.get("line_end"), int)
                        else None,
                        line_source=str(
                            block.metadata.get("line_match_strategy") or ""
                        ),
                    )
                )

        return tables

    def _extract_sections_from_toc(
        self,
        toc: list[dict],
        markdown: str,
    ) -> list[SectionAsset]:
        """Convert Marker TOC to SectionAsset list."""
        if not toc:
            return self.manifest_generator._parse_sections(markdown)

        sections = []
        existing_ids: set[str] = set()
        line_index = MarkdownLineSpanIndex(markdown)

        for idx, item in enumerate(toc, 1):
            title = item.get("title", "")
            base_id = f"sec_{idx}"
            sec_id = base_id
            if sec_id in existing_ids:
                counter = 2
                while f"{base_id}_{counter}" in existing_ids:
                    counter += 1
                sec_id = f"{base_id}_{counter}"
            section_span = line_index.find_section_span(
                title,
                page_hint=item.get("page", 0) or None,
            )
            sections.append(
                SectionAsset(
                    id=sec_id,
                    title=title,
                    level=item.get("level", 1),
                    page=item.get("page", 0),
                    start_line=section_span.start_line if section_span else 0,
                    end_line=section_span.end_line if section_span else 0,
                    preview=(
                        line_index.extract_preview(
                            section_span.start_line,
                            section_span.end_line,
                        )
                        if section_span is not None
                        else ""
                    ),
                )
            )
            existing_ids.add(sec_id)

        return sections

    async def _extract_and_save_images(
        self,
        doc_id: str,
        pdf_path: Path,
        *,
        page_map: list[int] | None = None,
    ) -> list[FigureAsset]:
        """Extract images from PDF with conservative filtering for textbook figures."""
        figures: list[FigureAsset] = []

        raw_images = self.pdf_extractor.extract_images(pdf_path)
        if not raw_images:
            return figures

        # Detect source from extractor type
        source = "pymupdf"
        if hasattr(self.pdf_extractor, "config"):
            source = "docling"

        # Extract figure captions for association
        page_captions: dict[int, list[dict]] = {}
        if hasattr(self.pdf_extractor, "extract_figure_captions"):
            page_captions = self.pdf_extractor.extract_figure_captions(pdf_path)

        min_px = self.profile.filters.min_figure_px
        if hasattr(self.pdf_extractor, "_MIN_FIGURE_PX"):
            min_px = int(getattr(self.pdf_extractor, "_MIN_FIGURE_PX", min_px))

        candidates_by_page: dict[int, list[dict[str, Any]]] = {}
        seen_hashes_by_page: dict[int, set[str]] = {}

        for img_data in raw_images:
            display_bytes = img_data.get("page_image_bytes") or img_data["image_bytes"]
            w = int(img_data.get("page_crop_width") or img_data["width"])
            h = int(img_data.get("page_crop_height") or img_data["height"])

            if w < min_px or h < min_px:
                continue

            longest_edge = max(w, h)
            shortest_edge = max(min(w, h), 1)
            aspect_ratio = longest_edge / shortest_edge
            if aspect_ratio >= 8.0:
                continue

            variance = self._image_variance(display_bytes)
            if variance < 5.0:
                continue

            page_num = int(img_data["page"])
            image_hash = hashlib.sha256(display_bytes).hexdigest()
            page_hashes = seen_hashes_by_page.setdefault(page_num, set())
            if image_hash in page_hashes:
                continue
            page_hashes.add(image_hash)

            area = w * h
            score = float(area) + (variance * 1500.0)
            if variance >= 15.0:
                score += 50_000.0
            if area >= 40_000:
                score += 25_000.0

            candidates_by_page.setdefault(page_num, []).append(
                {
                    **img_data,
                    "variance": variance,
                    "area": area,
                    "score": score,
                }
            )

        # Track which captions have been used (per page)
        used_captions: dict[int, set[int]] = {}

        for page_num in sorted(candidates_by_page):
            candidates = sorted(
                candidates_by_page[page_num],
                key=lambda item: (
                    float(item["score"]),
                    int(item["area"]),
                    float(item["variance"]),
                ),
                reverse=True,
            )
            caption_count = len(page_captions.get(page_num, []))
            if caption_count > 0:
                high_quality = [
                    item
                    for item in candidates
                    if float(item["variance"]) >= 15.0 or int(item["area"]) >= 30_000
                ]
                candidates = high_quality or candidates[:1]
                candidates = candidates[: max(caption_count * 4, 4)]
            else:
                candidates = [
                    item
                    for item in candidates
                    if float(item["variance"]) >= 15.0 and int(item["area"]) >= 40_000
                ][:2]
            candidates = sorted(candidates, key=self._figure_candidate_order_key)
            caps = page_captions.get(page_num, [])
            grouped_candidates, consumed_candidate_indexes, grouped_caption_indexes = (
                self._build_caption_group_candidates(
                    pdf_path,
                    page_num,
                    candidates,
                    caps,
                )
            )
            if grouped_candidates:
                candidates = grouped_candidates + [
                    item
                    for idx, item in enumerate(candidates)
                    if idx not in consumed_candidate_indexes
                ]
                used_captions.setdefault(page_num, set()).update(
                    grouped_caption_indexes
                )

            for local_index, img_data in enumerate(candidates, start=1):
                display_bytes = (
                    img_data.get("page_image_bytes") or img_data["image_bytes"]
                )
                display_ext = str(img_data.get("page_image_ext") or img_data["ext"])
                w = int(img_data.get("page_crop_width") or img_data["width"])
                h = int(img_data.get("page_crop_height") or img_data["height"])
                original_page = remap_page_number(page_num, page_map)

                # Generate figure ID using the curated local page order.
                fig_id = f"fig_{original_page}_{local_index}"

                raw_path = ""
                if self._should_save_raw_image(img_data, display_bytes, w, h):
                    raw_image_path = self.repository.save_image(
                        doc_id=doc_id,
                        image_id=f"{fig_id}_raw",
                        data=img_data["image_bytes"],
                        ext=img_data["ext"],
                    )
                    raw_path = str(raw_image_path)

                # Associate caption by geometry first; fall back to legacy FIFO.
                caption = img_data.get("caption", "")
                caption_bbox: list[float] = self._bbox_list(
                    img_data.get("caption_bbox")
                )
                caption_confidence = float(img_data.get("caption_confidence") or 0.0)
                if not caption:
                    caps = page_captions.get(page_num, [])
                    if page_num not in used_captions:
                        used_captions[page_num] = set()
                    match = self._match_caption_for_image(
                        img_data,
                        caps,
                        used_caption_indexes=used_captions[page_num],
                    )
                    if match is not None:
                        caption = str(match["caption"])
                        caption_bbox = list(match.get("bbox") or [])
                        caption_confidence = float(match.get("confidence") or 0.0)
                        used_captions[page_num].add(int(match["index"]))

                figure_bbox = self._bbox_list(img_data.get("bbox"))
                crop_bbox = self._bbox_list(img_data.get("page_crop_bbox"))
                extraction_strategy = str(img_data.get("extraction_strategy") or "")
                anchored_crop = self._render_caption_anchored_figure_crop(
                    pdf_path,
                    page_number=page_num,
                    image_data=img_data,
                    caption_bbox=caption_bbox,
                )
                if anchored_crop:
                    display_bytes = anchored_crop["image"]
                    display_ext = anchored_crop["ext"]
                    w = int(anchored_crop["width"])
                    h = int(anchored_crop["height"])
                    crop_bbox = list(anchored_crop["bbox"])
                    figure_bbox = list(anchored_crop["figure_bbox"])
                    extraction_strategy = "caption_anchor_page_crop"

                # Save display image. For PyMuPDF xobjects this is the page-region
                # crop, preserving PDF text/vector overlays that raw image objects miss.
                image_path = self.repository.save_image(
                    doc_id=doc_id,
                    image_id=fig_id,
                    data=display_bytes,
                    ext=display_ext,
                )

                figures.append(
                    FigureAsset(
                        id=fig_id,
                        page=original_page,
                        path=str(image_path),
                        ext=display_ext,
                        width=w,
                        height=h,
                        caption=caption,
                        raw_path=raw_path,
                        figure_bbox=figure_bbox,
                        crop_bbox=crop_bbox,
                        caption_bbox=caption_bbox,
                        caption_confidence=caption_confidence,
                        extraction_strategy=extraction_strategy,
                        figure_type="",
                        source=source,
                    )
                )

        return figures

    def _build_caption_group_candidates(
        self,
        pdf_path: Path,
        page_num: int,
        candidates: list[dict[str, Any]],
        captions: list[dict],
    ) -> tuple[list[dict[str, Any]], set[int], set[int]]:
        """Build one display crop per caption when a figure spans many xobjects."""
        grouped: list[dict[str, Any]] = []
        consumed_candidate_indexes: set[int] = set()
        consumed_caption_indexes: set[int] = set()

        caption_items: list[tuple[int, dict[str, Any], list[float]]] = []
        for caption_index, caption in enumerate(captions):
            caption_bbox = self._bbox_list(caption.get("bbox"))
            if caption_bbox:
                caption_items.append((caption_index, caption, caption_bbox))

        caption_items.sort(key=lambda item: (item[2][1], item[2][0]))

        for item_index, (caption_index, caption, caption_bbox) in enumerate(
            caption_items
        ):
            next_caption_bbox = (
                caption_items[item_index + 1][2]
                if item_index + 1 < len(caption_items)
                else None
            )
            matches: list[tuple[int, dict[str, Any]]] = []
            for candidate_index, candidate in enumerate(candidates):
                if candidate_index in consumed_candidate_indexes:
                    continue
                if self._candidate_belongs_to_caption(
                    candidate,
                    caption_bbox,
                    next_caption_bbox=next_caption_bbox,
                    caption_text=str(caption.get("caption") or ""),
                ):
                    matches.append((candidate_index, candidate))

            if len(matches) < 2:
                continue

            figure_bbox = self._union_bboxes(
                [
                    self._bbox_list(item.get("bbox") or item.get("page_crop_bbox"))
                    for _, item in matches
                ]
            )
            crop_bbox = self._union_bboxes([figure_bbox, caption_bbox])
            first = matches[0][1]
            rendered = self._render_caption_anchored_figure_crop(
                pdf_path,
                page_number=page_num,
                image_data=first,
                caption_bbox=caption_bbox,
            )
            extraction_strategy = (
                "caption_anchor_page_crop" if rendered else "caption_group_page_crop"
            )
            if not rendered:
                rendered = self._render_page_region_image(
                    pdf_path,
                    page_number=page_num,
                    bbox=crop_bbox,
                    padding=8.0,
                )
            if not rendered:
                continue

            grouped.append(
                {
                    **first,
                    "bbox": list(rendered.get("figure_bbox") or figure_bbox),
                    "page_image_bytes": rendered["image"],
                    "page_image_ext": rendered["ext"],
                    "page_crop_bbox": rendered["bbox"],
                    "page_crop_width": rendered["width"],
                    "page_crop_height": rendered["height"],
                    "caption": str(caption.get("caption") or ""),
                    "caption_bbox": caption_bbox,
                    "caption_confidence": 1.0,
                    "extraction_strategy": extraction_strategy,
                    "grouped_candidate_count": len(matches),
                    "save_raw_image": False,
                }
            )
            consumed_candidate_indexes.update(index for index, _ in matches)
            consumed_caption_indexes.add(caption_index)

        return grouped, consumed_candidate_indexes, consumed_caption_indexes

    @classmethod
    def _candidate_belongs_to_caption(
        cls,
        candidate: dict[str, Any],
        caption_bbox: list[float],
        *,
        next_caption_bbox: list[float] | None = None,
        caption_text: str = "",
    ) -> bool:
        """Return whether an image candidate is spatially part of a captioned figure."""
        image_bbox = cls._bbox_list(
            candidate.get("bbox") or candidate.get("page_crop_bbox")
        )
        if not image_bbox or not caption_bbox:
            return False
        img_x0, img_y0, img_x1, img_y1 = image_bbox
        cap_x0, cap_y0, cap_x1, cap_y1 = caption_bbox

        if next_caption_bbox:
            next_y0 = next_caption_bbox[1]
            image_center_y = (img_y0 + img_y1) / 2.0
            if img_y0 >= next_y0 - 8.0 or image_center_y >= next_y0 - 8.0:
                return False

        if img_y0 > cap_y0 + 8.0:
            return False
        vertical_gap = cls._vertical_gap(img_y0, img_y1, cap_y0, cap_y1)
        is_multipanel_caption = bool(
            re.search(r"\([A-Z]\)", caption_text)
            or re.search(r"\b[A-Z]\s*(?:and|,)\s*[A-Z]\b", caption_text)
        )
        max_vertical_gap = 260.0 if is_multipanel_caption else 150.0
        if vertical_gap > max_vertical_gap:
            return False

        expanded_cap_x0 = cap_x0 - 80.0
        expanded_cap_x1 = cap_x1 + 80.0
        horizontal_overlap = max(
            0.0,
            min(img_x1, expanded_cap_x1) - max(img_x0, expanded_cap_x0),
        )
        if horizontal_overlap > 0:
            return True
        image_center_x = (img_x0 + img_x1) / 2.0
        return expanded_cap_x0 <= image_center_x <= expanded_cap_x1

    @staticmethod
    def _union_bboxes(bboxes: list[list[float]]) -> list[float]:
        """Return a bbox enclosing all non-empty bboxes."""
        valid = [bbox for bbox in bboxes if len(bbox) >= 4]
        if not valid:
            return []
        return [
            round(min(bbox[0] for bbox in valid), 3),
            round(min(bbox[1] for bbox in valid), 3),
            round(max(bbox[2] for bbox in valid), 3),
            round(max(bbox[3] for bbox in valid), 3),
        ]

    @staticmethod
    def _render_page_region_image(
        pdf_path: Path,
        *,
        page_number: int,
        bbox: list[float],
        padding: float = 8.0,
        zoom: float = 2.0,
    ) -> dict[str, Any] | None:
        """Render a page-region crop for a grouped captioned figure."""
        if len(bbox) < 4 or page_number < 1:
            return None

        try:
            import fitz
        except Exception:
            return None

        try:
            with fitz.open(str(pdf_path)) as doc:
                if page_number > len(doc):
                    return None
                page = doc[page_number - 1]
                clip = fitz.Rect(
                    bbox[0] - padding,
                    bbox[1] - padding,
                    bbox[2] + padding,
                    bbox[3] + padding,
                )
                clip = clip & page.rect
                if clip.is_empty:
                    return None
                pix = page.get_pixmap(
                    matrix=fitz.Matrix(zoom, zoom),
                    clip=clip,
                    alpha=False,
                )
                return {
                    "image": cast("bytes", pix.tobytes("png")),
                    "ext": "png",
                    "width": pix.width,
                    "height": pix.height,
                    "bbox": [
                        round(float(clip.x0), 3),
                        round(float(clip.y0), 3),
                        round(float(clip.x1), 3),
                        round(float(clip.y1), 3),
                    ],
                }
        except Exception:
            return None

    def _render_caption_anchored_figure_crop(
        self,
        pdf_path: Path,
        *,
        page_number: int,
        image_data: dict[str, Any],
        caption_bbox: list[float],
    ) -> dict[str, Any] | None:
        """Render a wide figure crop when a tiny XObject belongs to a large captioned figure."""
        image_bbox = self._bbox_list(image_data.get("bbox"))
        if not image_bbox or not caption_bbox:
            return None
        if str(image_data.get("extraction_strategy") or "") != "xobject_page_crop":
            return None

        img_x0, img_y0, img_x1, img_y1 = image_bbox
        cap_x0, cap_y0, cap_x1, cap_y1 = caption_bbox
        image_width = max(img_x1 - img_x0, 1.0)
        image_height = max(img_y1 - img_y0, 1.0)
        caption_width = max(cap_x1 - cap_x0, 1.0)

        # Large textbook/vector figures can contain only one small embedded bitmap
        # (for example a probe illustration) while the rest is PDF vector/text.
        if caption_width < 300.0:
            return None
        if image_width >= caption_width * 0.35 or image_height >= 160.0:
            return None

        figure_height = min(380.0, max(220.0, caption_width * 0.70))
        crop_bbox = [
            cap_x0 - 8.0,
            max(0.0, cap_y0 - figure_height),
            cap_x1 + 8.0,
            cap_y1 + 4.0,
        ]
        rendered = self._render_page_region_image(
            pdf_path,
            page_number=page_number,
            bbox=crop_bbox,
            padding=0.0,
            zoom=2.5,
        )
        if not rendered:
            return None
        rendered["figure_bbox"] = [
            rendered["bbox"][0],
            rendered["bbox"][1],
            rendered["bbox"][2],
            round(float(cap_y0), 3),
        ]
        return rendered

    @staticmethod
    def _figure_candidate_order_key(item: dict[str, Any]) -> tuple[float, float, int]:
        """Return stable page reading order for figure IDs and captions."""
        bbox = DocumentService._bbox_list(
            item.get("bbox") or item.get("page_crop_bbox")
        )
        if bbox:
            return (bbox[1], bbox[0], int(item.get("index_on_page") or 0))
        return (0.0, 0.0, int(item.get("index_on_page") or 0))

    @staticmethod
    def _bbox_list(value: Any) -> list[float]:
        """Normalize optional bbox values to JSON-safe float lists."""
        if not isinstance(value, (list, tuple)) or len(value) < 4:
            return []
        return [round(float(item), 3) for item in value[:4]]

    def _should_save_raw_image(
        self,
        image_data: dict[str, Any],
        display_bytes: bytes,
        display_width: int,
        display_height: int,
    ) -> bool:
        """Return whether the original XObject is useful enough to persist."""
        if not image_data.get("page_image_bytes") or not image_data.get(
            "save_raw_image",
            True,
        ):
            return False

        raw_bytes = image_data.get("image_bytes")
        if not isinstance(raw_bytes, (bytes, bytearray)) or not raw_bytes:
            return False
        if hashlib.sha256(raw_bytes).digest() == hashlib.sha256(display_bytes).digest():
            return False

        raw_width = int(image_data.get("width") or 0)
        raw_height = int(image_data.get("height") or 0)
        if (
            raw_width < self._get_min_figure_px()
            or raw_height < self._get_min_figure_px()
        ):
            return False

        raw_area = max(raw_width * raw_height, 1)
        display_area = max(display_width * display_height, 1)
        if display_area <= raw_area * 1.10:
            return False

        return self._image_variance(bytes(raw_bytes)) >= 5.0

    @classmethod
    def _match_caption_for_image(
        cls,
        image_data: dict[str, Any],
        captions: list[dict],
        *,
        used_caption_indexes: set[int],
    ) -> dict[str, Any] | None:
        """Match the most likely caption to an image using page geometry."""
        figure_bbox = cls._bbox_list(image_data.get("bbox")) or cls._bbox_list(
            image_data.get("page_crop_bbox")
        )
        if not figure_bbox:
            return None

        fig_x0, fig_y0, fig_x1, fig_y1 = figure_bbox
        fig_center_x = (fig_x0 + fig_x1) / 2.0
        fig_width = max(fig_x1 - fig_x0, 1.0)
        best: dict[str, Any] | None = None
        best_score = -1.0

        for idx, cap in enumerate(captions):
            if idx in used_caption_indexes:
                continue
            caption = str(cap.get("caption") or "")
            bbox = cls._bbox_list(cap.get("bbox"))
            if not caption or not bbox:
                continue
            cap_x0, cap_y0, cap_x1, cap_y1 = bbox
            if cap_y0 < fig_y1 - 8.0:
                continue
            cap_center_x = (cap_x0 + cap_x1) / 2.0
            vertical_gap = cls._vertical_gap(fig_y0, fig_y1, cap_y0, cap_y1)
            if vertical_gap > 180.0:
                continue
            if any(
                other_idx != idx
                and other_idx not in used_caption_indexes
                and (other_bbox := cls._bbox_list(other.get("bbox")))
                and fig_y1 < other_bbox[1] < cap_y0
                for other_idx, other in enumerate(captions)
            ):
                continue
            horizontal_offset = abs(fig_center_x - cap_center_x)
            horizontal_overlap = max(0.0, min(fig_x1, cap_x1) - max(fig_x0, cap_x0))
            overlap_ratio = horizontal_overlap / min(
                fig_width,
                max(cap_x1 - cap_x0, 1.0),
            )

            below_bonus = 0.35 if cap_y0 >= fig_y0 else 0.0
            overlap_bonus = min(overlap_ratio, 1.0) * 0.35
            distance_penalty = min(vertical_gap / 240.0, 1.0) * 0.45
            offset_penalty = min(horizontal_offset / max(fig_width, 1.0), 1.0) * 0.30
            score = (
                0.75 + below_bonus + overlap_bonus - distance_penalty - offset_penalty
            )

            if score > best_score:
                best_score = score
                best = {
                    "index": idx,
                    "caption": caption,
                    "bbox": bbox,
                    "confidence": max(0.0, min(score, 1.0)),
                }

        if best is None or float(best["confidence"]) < 0.50:
            return None
        return best

    @staticmethod
    def _vertical_gap(
        fig_y0: float,
        fig_y1: float,
        cap_y0: float,
        cap_y1: float,
    ) -> float:
        """Return vertical gap between figure and caption boxes."""
        if cap_y0 >= fig_y1:
            return cap_y0 - fig_y1
        if fig_y0 >= cap_y1:
            return fig_y0 - cap_y1
        return 0.0

    async def _extract_tables(
        self,
        pdf_path: Path,
        *,
        page_map: list[int] | None = None,
    ) -> list[TableAsset]:
        """
        Extract tables from PDF.

        Supports:
        - PyMuPDF: find_tables() - heuristic, good for simple grid tables
        - Docling (optional): TableFormer - AI-based, better for complex tables
        """
        # Check if extractor supports table extraction
        if not hasattr(self.pdf_extractor, "extract_tables"):
            return []  # Will fall back to markdown parsing

        # Detect source from extractor type
        source = "pymupdf"
        if hasattr(self.pdf_extractor, "config"):
            source = "docling"

        try:
            raw_tables = self.pdf_extractor.extract_tables(pdf_path)

            tables: list[TableAsset] = []
            for tab_data in raw_tables:
                page_number = remap_page_number(tab_data.get("page", 1), page_map)
                tables.append(
                    TableAsset(
                        id=tab_data.get("id", f"tab_{len(tables) + 1}"),
                        page=page_number,
                        caption=tab_data.get("caption", ""),
                        preview=tab_data.get("preview", ""),
                        markdown=tab_data.get("markdown", ""),
                        row_count=tab_data.get("row_count", 0),
                        col_count=tab_data.get("col_count", 0),
                        has_header=tab_data.get("has_header", True),
                        source=source,
                    )
                )
            return tables

        except Exception as e:
            import logging

            logging.warning(f"Table extraction failed: {e}")
            return []

    async def list_documents(self) -> list[DocumentSummary]:
        """List all processed documents."""
        return self.repository.list_documents()

    async def get_manifest(self, doc_id: str) -> DocumentManifest | None:
        """Get manifest for a specific document."""
        return self.repository.load_manifest(doc_id)

    async def document_exists(self, doc_id: str) -> bool:
        """Check if a document exists."""
        return self.repository.document_exists(doc_id)

    async def delete_document(self, doc_id: str) -> dict[str, Any]:
        """Delete a stored PDF document and its local artifacts."""
        manifest = self.repository.load_manifest(doc_id)
        if manifest is None:
            return {"success": False, "error": f"Document not found: {doc_id}"}

        deleted = self.repository.delete_document(doc_id)
        if not deleted:
            return {
                "success": False,
                "error": f"Failed to delete document directory for {doc_id}",
            }

        warnings: list[str] = []
        knowledge_graph_status: str | None = None
        if self.knowledge_graph and self.knowledge_graph.is_available:
            try:
                kg_result = await self.knowledge_graph.delete_document(doc_id)
                knowledge_graph_status = str(kg_result.get("status", "unknown"))
                if knowledge_graph_status not in {"success", "not_found"}:
                    warnings.append(
                        "Knowledge graph deletion did not complete successfully; local artifacts were deleted first."
                    )
            except Exception as exc:
                warnings.append(
                    "Knowledge graph deletion failed; local artifacts were deleted first. "
                    f"Reason: {exc}"
                )

        result = {
            "success": True,
            "doc_id": doc_id,
            "filename": manifest.filename,
            "warnings": warnings,
        }
        if knowledge_graph_status is not None:
            result["knowledge_graph_status"] = knowledge_graph_status
        return result

    async def convert_pdf_to_docx(
        self,
        doc_id: str,
        output_path: str | None = None,
        *,
        mode: str = "content",
    ) -> dict[str, Any]:
        """
        Convert an ingested PDF document to DOCX.

        Supported modes:
        - ``content``: rebuild a readable DOCX from extracted markdown and figures.
        - ``fidelity``: currently unsupported because PDF ETL is not layout-reversible.
        """
        if mode != "content":
            return {
                "success": False,
                "error": (
                    "PDF → DOCX currently supports content mode only. "
                    "Layout-fidelity reconstruction is not available."
                ),
            }

        manifest = self.repository.load_manifest(doc_id)
        if manifest is None:
            return {"success": False, "error": f"Document not found: {doc_id}"}

        markdown = self.repository.load_markdown(doc_id)
        if markdown is None:
            return {
                "success": False,
                "error": f"Markdown content not found for {doc_id}",
            }

        doc_dir = self.repository.get_doc_dir(doc_id)
        try:
            out_path = resolve_document_output_path(
                doc_dir,
                output_path,
                default_name="converted_from_pdf.docx",
                allowed_suffixes={".docx"},
            )
        except ValueError as e:
            return {"success": False, "error": str(e)}

        try:
            self._build_docx_from_markdown(markdown, manifest, out_path)
        except Exception as e:
            return {"success": False, "error": str(e)}

        return {
            "success": True,
            "doc_id": doc_id,
            "output_path": str(out_path),
            "mode": mode,
            "figures_embedded": len(manifest.assets.figures),
            "tables_found": len(manifest.assets.tables),
        }

    async def convert_pdf_to_pptx(
        self,
        doc_id: str,
        output_path: str | None = None,
        *,
        mode: str = "content",
    ) -> dict[str, Any]:
        """
        Convert an ingested PDF document to PPTX slides.

        Supported modes:
        - ``content``: slide-oriented rendering from extracted markdown + figures.
        """
        if mode != "content":
            return {
                "success": False,
                "error": (
                    "PDF → PPTX currently supports content mode only. "
                    "Layout-fidelity reconstruction is not available."
                ),
            }

        manifest = self.repository.load_manifest(doc_id)
        if manifest is None:
            return {"success": False, "error": f"Document not found: {doc_id}"}

        markdown = self.repository.load_markdown(doc_id)
        if markdown is None:
            return {
                "success": False,
                "error": f"Markdown content not found for {doc_id}",
            }

        doc_dir = self.repository.get_doc_dir(doc_id)
        try:
            out_path = resolve_document_output_path(
                doc_dir,
                output_path,
                default_name="converted_from_pdf.pptx",
                allowed_suffixes={".pptx"},
            )
        except ValueError as e:
            return {"success": False, "error": str(e)}

        try:
            build_stats = self._build_pptx_from_markdown(markdown, manifest, out_path)
        except Exception as e:
            return {"success": False, "error": str(e)}

        return {
            "success": True,
            "doc_id": doc_id,
            "output_path": str(out_path),
            "mode": mode,
            "slides_created": build_stats.get("total_slides", 0),
            "figure_slides": build_stats.get("figure_slides", 0),
        }

    def _build_docx_from_markdown(
        self,
        markdown: str,
        manifest: DocumentManifest,
        output_path: Path,
    ) -> None:
        """Render extracted markdown into a readable DOCX document."""
        from docx import Document
        from docx.enum.text import WD_BREAK
        from docx.image.exceptions import UnrecognizedImageError
        from docx.shared import Cm

        document = Document()
        if manifest.title:
            document.add_heading(manifest.title, level=0)
            document.core_properties.title = manifest.title

        lines = markdown.splitlines()
        index = 0
        while index < len(lines):
            raw_line = lines[index].rstrip()
            stripped = raw_line.strip()

            if not stripped:
                index += 1
                continue

            if stripped.startswith("<!--") and stripped.endswith("-->"):
                if stripped.lower().startswith("<!-- page") and document.paragraphs:
                    document.paragraphs[-1].add_run().add_break(WD_BREAK.PAGE)
                index += 1
                continue

            if stripped.startswith("#"):
                level = min(len(stripped) - len(stripped.lstrip("#")), 9)
                document.add_heading(stripped[level:].strip(), level=level)
                index += 1
                continue

            if self._is_table_start(lines, index):
                index = self._append_markdown_table(document, lines, index)
                continue

            if self._is_list_item(stripped):
                index = self._append_markdown_list(document, lines, index)
                continue

            paragraph_lines = [stripped]
            index += 1
            while index < len(lines):
                next_line = lines[index].strip()
                if (
                    not next_line
                    or next_line.startswith(("<!--", "#"))
                    or self._is_list_item(next_line)
                    or self._is_table_start(lines, index)
                ):
                    break
                paragraph_lines.append(next_line)
                index += 1

            document.add_paragraph(" ".join(paragraph_lines))

        if manifest.assets.figures:
            document.add_page_break()
            document.add_heading("Extracted Figures", level=1)
            with tempfile.TemporaryDirectory() as image_tmp_dir:
                for figure in manifest.assets.figures:
                    if figure.caption:
                        document.add_paragraph(figure.caption)
                    figure_path = Path(figure.path)
                    if figure_path.exists():
                        try:
                            document.add_picture(str(figure_path), width=Cm(15))
                        except UnrecognizedImageError:
                            compatible_path = self._normalize_image_for_docx(
                                figure_path,
                                Path(image_tmp_dir),
                            )
                            document.add_picture(str(compatible_path), width=Cm(15))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(str(output_path))

    @staticmethod
    def _normalize_image_for_docx(image_path: Path, output_dir: Path) -> Path:
        """Re-encode readable images into a header form python-docx accepts."""
        from PIL import Image

        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{image_path.stem}.png"
        with Image.open(image_path) as image:
            if image.mode not in {"RGB", "RGBA"}:
                converted = image.convert("RGB")
                converted.save(output_path, format="PNG")
            else:
                image.save(output_path, format="PNG")
        return output_path

    def _build_pptx_from_markdown(
        self,
        markdown: str,
        manifest: DocumentManifest,
        output_path: Path,
    ) -> dict[str, int]:
        """Render extracted markdown into a slide deck."""
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        title_or_fallback = manifest.title or manifest.filename
        title_layout = self._get_slide_layout(presentation, 0)
        content_layout = self._get_slide_layout(presentation, 1, default_index=0)
        figure_layout = self._get_slide_layout(
            presentation, 5, default_index=content_layout
        )
        blank_layout = self._get_slide_layout(
            presentation, 6, default_index=content_layout
        )

        slides = self._segment_markdown_to_slides(markdown, title_or_fallback)
        for slide_title, items in slides:
            layout = content_layout if items else title_layout
            slide = presentation.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
            if items and len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                for index, (text, level) in enumerate(items):
                    paragraph = (
                        text_frame.paragraphs[0]
                        if index == 0 and text_frame.paragraphs
                        else text_frame.add_paragraph()
                    )
                    paragraph.text = text
                    paragraph.level = max(level, 0)

        figure_slides = 0
        for figure_index, figure in enumerate(manifest.assets.figures, start=1):
            figure_path = Path(figure.path)
            if not figure_path.exists():
                continue

            slide = presentation.slides.add_slide(figure_layout)
            if slide.shapes.title:
                slide.shapes.title.text = (
                    figure.caption or f"Figure {figure_index}" or title_or_fallback
                )

            left = Inches(0.75)
            top = Inches(1.5)
            max_width = Inches(9)
            try:
                slide.shapes.add_picture(
                    str(figure_path), left=left, top=top, width=max_width
                )
            except Exception:
                slide = presentation.slides.add_slide(blank_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = (
                        figure.caption or f"Figure {figure_index}" or title_or_fallback
                    )
                slide.shapes.add_picture(
                    str(figure_path), left=left, top=top, width=max_width
                )

            figure_slides += 1

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))

        return {
            "total_slides": len(presentation.slides),
            "figure_slides": figure_slides,
        }

    @staticmethod
    def _is_list_item(line: str) -> bool:
        return bool(re.match(r"^([-*+]\s+|\d+[.)]\s+)", line))

    @staticmethod
    def _is_table_start(lines: list[str], index: int) -> bool:
        if index + 1 >= len(lines):
            return False
        header = lines[index].strip()
        separator = lines[index + 1].strip()
        return (
            header.startswith("|")
            and header.endswith("|")
            and separator.startswith("|")
            and separator.endswith("|")
            and set(separator.replace("|", "").replace(" ", "")) <= {"-", ":"}
        )

    @staticmethod
    def _split_markdown_row(line: str) -> list[str]:
        return [cell.strip() for cell in line.strip().strip("|").split("|")]

    def _append_markdown_table(
        self, document: Any, lines: list[str], index: int
    ) -> int:
        header_cells = self._split_markdown_row(lines[index])
        index += 2  # Skip header + separator

        rows: list[list[str]] = []
        while index < len(lines):
            current = lines[index].strip()
            if not (current.startswith("|") and current.endswith("|")):
                break
            rows.append(self._split_markdown_row(current))
            index += 1

        table = document.add_table(rows=1, cols=len(header_cells))
        table.style = "Table Grid"
        for col_index, value in enumerate(header_cells):
            table.rows[0].cells[col_index].text = value

        for row in rows:
            cells = table.add_row().cells
            for col_index, value in enumerate(row[: len(header_cells)]):
                cells[col_index].text = value

        return index

    def _append_markdown_list(self, document: Any, lines: list[str], index: int) -> int:
        while index < len(lines):
            stripped = lines[index].strip()
            if not self._is_list_item(stripped):
                break
            style = (
                "List Number" if re.match(r"^\d+[.)]\s+", stripped) else "List Bullet"
            )
            text = re.sub(r"^([-*+]\s+|\d+[.)]\s+)", "", stripped)
            document.add_paragraph(text, style=style)
            index += 1
        return index

    @staticmethod
    def _get_slide_layout(
        presentation: Any,
        index: int,
        *,
        default_index: int | Any = 0,
    ) -> Any:
        """Safely fetch a slide layout, falling back when missing."""
        layouts = presentation.slide_layouts
        if 0 <= index < len(layouts):
            return layouts[index]
        if isinstance(default_index, int):
            safe_index = min(max(default_index, 0), len(layouts) - 1)
            return layouts[safe_index]
        return default_index

    def _segment_markdown_to_slides(
        self,
        markdown: str,
        fallback_title: str,
    ) -> list[tuple[str, list[tuple[str, int]]]]:
        """
        Partition markdown into slide-sized chunks anchored by headings.

        Returns:
            List of (title, [(text, level), ...]) tuples.
        """
        slides: list[tuple[str, list[tuple[str, int]]]] = []
        lines = markdown.splitlines()
        index = 0

        current_title = fallback_title
        current_items: list[tuple[str, int]] = []

        while index < len(lines):
            raw_line = lines[index]
            stripped = raw_line.strip()
            if not stripped:
                index += 1
                continue

            if stripped.startswith("<!--"):
                index += 1
                continue

            if stripped.startswith("#"):
                if current_items or slides:
                    slides.append((current_title, current_items))
                level = len(stripped) - len(stripped.lstrip("#"))
                heading_text = stripped[level:].strip() or fallback_title
                current_title = heading_text
                current_items = []
                index += 1
                continue

            if self._is_table_start(lines, index):
                header_cells = self._split_markdown_row(lines[index])
                index += 2  # Skip header and separator
                preview_row: list[str] = []
                while index < len(lines):
                    row_line = lines[index].strip()
                    if not (row_line.startswith("|") and row_line.endswith("|")):
                        break
                    preview_row = self._split_markdown_row(row_line)
                    index += 1
                    if preview_row:
                        break

                preview = " | ".join(header_cells)
                if preview_row:
                    preview += " — " + " | ".join(preview_row)
                current_items.append((f"Table: {preview}", 0))
                continue

            if self._is_list_item(stripped):
                indent = len(raw_line) - len(raw_line.lstrip(" "))
                level = min(indent // 2, 4)
                text = re.sub(r"^([-*+]\s+|\d+[.)]\s+)", "", stripped)
                current_items.append((text, level))
                index += 1
                continue

            paragraph_lines = [stripped]
            index += 1
            while index < len(lines):
                next_line = lines[index].strip()
                if (
                    not next_line
                    or next_line.startswith(("<!--", "#"))
                    or self._is_list_item(next_line)
                    or self._is_table_start(lines, index)
                ):
                    break
                paragraph_lines.append(next_line)
                index += 1

            merged_paragraph = " ".join(paragraph_lines)
            current_items.append((merged_paragraph, 0))

        if current_items or not slides:
            slides.append((current_title, current_items))

        return slides
