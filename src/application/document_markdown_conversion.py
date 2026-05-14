"""Markdown-to-office renderers used by DocumentService conversions."""

from __future__ import annotations

import re
import tempfile
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from src.domain.entities import DocumentManifest


class MarkdownConversionMixin:
    """Private helpers for rendering extracted Markdown into DOCX/PPTX."""

    def _build_docx_from_markdown(
        self,
        markdown: str,
        manifest: DocumentManifest,
        output_path: Path,
    ) -> None:
        """Render extracted markdown into a readable DOCX document."""
        from docx import Document
        from docx.enum.text import WD_BREAK
        from docx.image.exceptions import UnrecognizedImageError
        from docx.shared import Cm

        document = Document()
        if manifest.title:
            document.add_heading(manifest.title, level=0)
            document.core_properties.title = manifest.title

        lines = markdown.splitlines()
        index = 0
        while index < len(lines):
            raw_line = lines[index].rstrip()
            stripped = raw_line.strip()

            if not stripped:
                index += 1
                continue

            if stripped.startswith("<!--") and stripped.endswith("-->"):
                if stripped.lower().startswith("<!-- page") and document.paragraphs:
                    document.paragraphs[-1].add_run().add_break(WD_BREAK.PAGE)
                index += 1
                continue

            if stripped.startswith("#"):
                level = min(len(stripped) - len(stripped.lstrip("#")), 9)
                document.add_heading(stripped[level:].strip(), level=level)
                index += 1
                continue

            if self._is_table_start(lines, index):
                index = self._append_markdown_table(document, lines, index)
                continue

            if self._is_list_item(stripped):
                index = self._append_markdown_list(document, lines, index)
                continue

            paragraph_lines = [stripped]
            index += 1
            while index < len(lines):
                next_line = lines[index].strip()
                if (
                    not next_line
                    or next_line.startswith(("<!--", "#"))
                    or self._is_list_item(next_line)
                    or self._is_table_start(lines, index)
                ):
                    break
                paragraph_lines.append(next_line)
                index += 1

            document.add_paragraph(" ".join(paragraph_lines))

        if manifest.assets.figures:
            document.add_page_break()
            document.add_heading("Extracted Figures", level=1)
            with tempfile.TemporaryDirectory() as image_tmp_dir:
                for figure in manifest.assets.figures:
                    if figure.caption:
                        document.add_paragraph(figure.caption)
                    figure_path = Path(figure.path)
                    if figure_path.exists():
                        try:
                            document.add_picture(str(figure_path), width=Cm(15))
                        except UnrecognizedImageError:
                            compatible_path = self._normalize_image_for_docx(
                                figure_path,
                                Path(image_tmp_dir),
                            )
                            document.add_picture(str(compatible_path), width=Cm(15))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(str(output_path))

    @staticmethod
    def _normalize_image_for_docx(image_path: Path, output_dir: Path) -> Path:
        """Re-encode readable images into a header form python-docx accepts."""
        from PIL import Image

        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{image_path.stem}.png"
        with Image.open(image_path) as image:
            if image.mode not in {"RGB", "RGBA"}:
                converted = image.convert("RGB")
                converted.save(output_path, format="PNG")
            else:
                image.save(output_path, format="PNG")
        return output_path

    def _build_pptx_from_markdown(
        self,
        markdown: str,
        manifest: DocumentManifest,
        output_path: Path,
    ) -> dict[str, int]:
        """Render extracted markdown into a slide deck."""
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        title_or_fallback = manifest.title or manifest.filename
        title_layout = self._get_slide_layout(presentation, 0)
        content_layout = self._get_slide_layout(presentation, 1, default_index=0)
        figure_layout = self._get_slide_layout(
            presentation, 5, default_index=content_layout
        )
        blank_layout = self._get_slide_layout(
            presentation, 6, default_index=content_layout
        )

        slides = self._segment_markdown_to_slides(markdown, title_or_fallback)
        for slide_title, items in slides:
            layout = content_layout if items else title_layout
            slide = presentation.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
            if items and len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                for index, (text, level) in enumerate(items):
                    paragraph = (
                        text_frame.paragraphs[0]
                        if index == 0 and text_frame.paragraphs
                        else text_frame.add_paragraph()
                    )
                    paragraph.text = text
                    paragraph.level = max(level, 0)

        figure_slides = 0
        for figure_index, figure in enumerate(manifest.assets.figures, start=1):
            figure_path = Path(figure.path)
            if not figure_path.exists():
                continue

            slide = presentation.slides.add_slide(figure_layout)
            if slide.shapes.title:
                slide.shapes.title.text = (
                    figure.caption or f"Figure {figure_index}" or title_or_fallback
                )

            left = Inches(0.75)
            top = Inches(1.5)
            max_width = Inches(9)
            try:
                slide.shapes.add_picture(
                    str(figure_path), left=left, top=top, width=max_width
                )
            except Exception:
                slide = presentation.slides.add_slide(blank_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = (
                        figure.caption or f"Figure {figure_index}" or title_or_fallback
                    )
                slide.shapes.add_picture(
                    str(figure_path), left=left, top=top, width=max_width
                )

            figure_slides += 1

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))

        return {
            "total_slides": len(presentation.slides),
            "figure_slides": figure_slides,
        }

    @staticmethod
    def _is_list_item(line: str) -> bool:
        return bool(re.match(r"^([-*+]\s+|\d+[.)]\s+)", line))

    @staticmethod
    def _is_table_start(lines: list[str], index: int) -> bool:
        if index + 1 >= len(lines):
            return False
        header = lines[index].strip()
        separator = lines[index + 1].strip()
        return (
            header.startswith("|")
            and header.endswith("|")
            and separator.startswith("|")
            and separator.endswith("|")
            and set(separator.replace("|", "").replace(" ", "")) <= {"-", ":"}
        )

    @staticmethod
    def _split_markdown_row(line: str) -> list[str]:
        return [cell.strip() for cell in line.strip().strip("|").split("|")]

    def _append_markdown_table(
        self, document: Any, lines: list[str], index: int
    ) -> int:
        header_cells = self._split_markdown_row(lines[index])
        index += 2

        rows: list[list[str]] = []
        while index < len(lines):
            current = lines[index].strip()
            if not (current.startswith("|") and current.endswith("|")):
                break
            rows.append(self._split_markdown_row(current))
            index += 1

        table = document.add_table(rows=1, cols=len(header_cells))
        table.style = "Table Grid"
        for col_index, value in enumerate(header_cells):
            table.rows[0].cells[col_index].text = value

        for row in rows:
            cells = table.add_row().cells
            for col_index, value in enumerate(row[: len(header_cells)]):
                cells[col_index].text = value

        return index

    def _append_markdown_list(self, document: Any, lines: list[str], index: int) -> int:
        while index < len(lines):
            stripped = lines[index].strip()
            if not self._is_list_item(stripped):
                break
            style = (
                "List Number" if re.match(r"^\d+[.)]\s+", stripped) else "List Bullet"
            )
            text = re.sub(r"^([-*+]\s+|\d+[.)]\s+)", "", stripped)
            document.add_paragraph(text, style=style)
            index += 1
        return index

    @staticmethod
    def _get_slide_layout(
        presentation: Any,
        index: int,
        *,
        default_index: int | Any = 0,
    ) -> Any:
        """Safely fetch a slide layout, falling back when missing."""
        layouts = presentation.slide_layouts
        if 0 <= index < len(layouts):
            return layouts[index]
        if isinstance(default_index, int):
            safe_index = min(max(default_index, 0), len(layouts) - 1)
            return layouts[safe_index]
        return default_index

    def _segment_markdown_to_slides(
        self,
        markdown: str,
        fallback_title: str,
    ) -> list[tuple[str, list[tuple[str, int]]]]:
        """
        Partition markdown into slide-sized chunks anchored by headings.

        Returns:
            List of (title, [(text, level), ...]) tuples.
        """
        slides: list[tuple[str, list[tuple[str, int]]]] = []
        lines = markdown.splitlines()
        index = 0

        current_title = fallback_title
        current_items: list[tuple[str, int]] = []

        while index < len(lines):
            raw_line = lines[index]
            stripped = raw_line.strip()
            if not stripped:
                index += 1
                continue

            if stripped.startswith("<!--"):
                index += 1
                continue

            if stripped.startswith("#"):
                if current_items or slides:
                    slides.append((current_title, current_items))
                level = len(stripped) - len(stripped.lstrip("#"))
                heading_text = stripped[level:].strip() or fallback_title
                current_title = heading_text
                current_items = []
                index += 1
                continue

            if self._is_table_start(lines, index):
                header_cells = self._split_markdown_row(lines[index])
                index += 2
                preview_row: list[str] = []
                while index < len(lines):
                    row_line = lines[index].strip()
                    if not (row_line.startswith("|") and row_line.endswith("|")):
                        break
                    preview_row = self._split_markdown_row(row_line)
                    index += 1
                    if preview_row:
                        break

                preview = " | ".join(header_cells)
                if preview_row:
                    preview += " — " + " | ".join(preview_row)
                current_items.append((f"Table: {preview}", 0))
                continue

            if self._is_list_item(stripped):
                indent = len(raw_line) - len(raw_line.lstrip(" "))
                level = min(indent // 2, 4)
                text = re.sub(r"^([-*+]\s+|\d+[.)]\s+)", "", stripped)
                current_items.append((text, level))
                index += 1
                continue

            paragraph_lines = [stripped]
            index += 1
            while index < len(lines):
                next_line = lines[index].strip()
                if (
                    not next_line
                    or next_line.startswith(("<!--", "#"))
                    or self._is_list_item(next_line)
                    or self._is_table_start(lines, index)
                ):
                    break
                paragraph_lines.append(next_line)
                index += 1

            merged_paragraph = " ".join(paragraph_lines)
            current_items.append((merged_paragraph, 0))

        if current_items or not slides:
            slides.append((current_title, current_items))

        return slides
